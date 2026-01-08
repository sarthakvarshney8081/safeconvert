from fastapi import APIRouter, UploadFile, File
from core.processor import save_upload_file, cleanup_file

from fastapi import APIRouter, UploadFile, File, HTTPException, BackgroundTasks, Response
from fastapi.responses import FileResponse, StreamingResponse
from core.processor import save_upload_file, cleanup_file, UPLOAD_DIR
import os
import uuid
import img2pdf
from pdf2image import convert_from_path
import subprocess
import zipfile
import shutil
import io

router = APIRouter(tags=["convert"])

@router.post("/office-to-pdf")
async def convert_office(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    temp_file = None
    output_dir = os.path.join(UPLOAD_DIR, f"office_{uuid.uuid4()}")
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        temp_file = await save_upload_file(file)
        
        # Check if LibreOffice is available (fallback? error?)
        # Command: libreoffice --headless --convert-to pdf <file> --outdir <dir>
        process = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', temp_file, '--outdir', output_dir],
            capture_output=True,
            text=True
        )
        
        if process.returncode != 0:
             raise Exception(f"LibreOffice conversion failed: {process.stderr}")

        # Find the output pdf
        files_in_output = os.listdir(output_dir)
        pdf_files = [f for f in files_in_output if f.lower().endswith('.pdf')]
        
        if not pdf_files:
             raise Exception("Output PDF not found after conversion")
             
        output_pdf_path = os.path.join(output_dir, pdf_files[0])
        
        # Cleanup input immediately
        if temp_file: cleanup_file(temp_file)
        
        # Schedule directory cleanup? 
        # Background task can remove the whole folder after serving
        background_tasks.add_task(shutil.rmtree, output_dir, ignore_errors=True) # Recursive delete

        return FileResponse(output_pdf_path, filename=pdf_files[0], media_type="application/pdf")

    except Exception as e:
        if temp_file: cleanup_file(temp_file)
        shutil.rmtree(output_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

@router.post("/image-to-pdf")
async def convert_image_to_pdf(background_tasks: BackgroundTasks, files: list[UploadFile] = File(...)):
    temp_files = []
    try:
        # Save images
        image_paths = []
        for file in files:
            path = await save_upload_file(file)
            temp_files.append(path)
            image_paths.append(path)

        output_filename = f"img2pdf_{uuid.uuid4()}.pdf"
        output_path = os.path.join(UPLOAD_DIR, output_filename)
        
        # Convert
        # img2pdf requires direct image bytes or paths.
        # But it fails on RGBA. We should convert to RGB via Pillow if needed.
        from PIL import Image
        
        valid_pdfs = []
        for path in image_paths:
             try:
                 img = Image.open(path)
                 if img.mode == 'RGBA':
                     # Convert to RGB
                     rgb_path = path + "_rgb.jpg"
                     img.convert('RGB').save(rgb_path, "JPEG")
                     valid_pdfs.append(rgb_path)
                     background_tasks.add_task(cleanup_file, rgb_path)
                 else:
                     valid_pdfs.append(path)
             except Exception:
                 # If not an image PIL can read, skip or try direct
                 valid_pdfs.append(path)
                 
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(valid_pdfs))

        for path in temp_files:
            background_tasks.add_task(cleanup_file, path)

        return FileResponse(output_path, filename="converted.pdf", media_type="application/pdf")

    except Exception as e:
        for path in temp_files:
            cleanup_file(path)
        raise HTTPException(status_code=500, detail=f"Image to PDF failed: {str(e)}")

@router.post("/pdf-to-image")
async def convert_pdf_to_image(background_tasks: BackgroundTasks, file: UploadFile = File(...), fmt: str = "png"):
    temp_file = None
    output_dir = os.path.join(UPLOAD_DIR, f"pdf2img_{uuid.uuid4()}")
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        temp_file = await save_upload_file(file)
        
        images = convert_from_path(temp_file)
        print(f"DEBUG: Found {len(images)} images in PDF", flush=True)
        
        # Always return a ZIP for consistency and to avoid extension issues on macOS
        zip_path = os.path.join(UPLOAD_DIR, f"images_{uuid.uuid4()}.zip")
        
        # Save images to folder first
        for i, image in enumerate(images):
            img_name = f"page_{i+1}.{fmt}"
            img_path = os.path.join(output_dir, img_name)
            if fmt.lower() == 'png' and image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(img_path, fmt.upper())

        # Use system ZIP (standard tool) for maximum macOS compatibility
        # -0: no compression (stored), -j: junk paths (no dirs)
        try:
            subprocess.run(
                ['zip', '-0', '-j', zip_path] + [os.path.join(output_dir, f) for f in os.listdir(output_dir)],
                check=True,
                capture_output=True
            )
        except subprocess.CalledProcessError as e:
            print(f"DEBUG: ZIP command failed: {e.stderr.decode()}", flush=True)
            raise Exception("System ZIP creation failed")
        
        # Log final file size for verification
        zip_size = os.path.getsize(zip_path)
        print(f"DEBUG: Final ZIP created at {zip_path}, size: {zip_size} bytes", flush=True)

        # Cleanup intermediate image files and input PDF
        background_tasks.add_task(shutil.rmtree, output_dir, ignore_errors=True)
        background_tasks.add_task(cleanup_file, temp_file)
        # Note: ZIP file will be cleaned up by the periodic task in main.py (10m TTL)

        return FileResponse(
            zip_path, 
            filename="converted_images.zip", 
            media_type="application/zip"
        )

    except Exception as e:
        if temp_file: cleanup_file(temp_file)
        shutil.rmtree(output_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"PDF to Image failed: {str(e)}")

@router.post("/pdf-to-excel")
async def convert_pdf_to_excel(
    background_tasks: BackgroundTasks, 
    file: UploadFile = File(...)
):
    import pdfplumber
    import pandas as pd
    
    temp_file = None
    output_filename = f"excel_{uuid.uuid4()}.xlsx"
    output_path = os.path.join(UPLOAD_DIR, output_filename)
    
    try:
        temp_file = await save_upload_file(file)
        
        # Strategy: Extract tables from all pages and merge into one sheet or multiple?
        # Let's merge into one big DataFrame for simplicity, or separate sheets.
        # "One sheet per page" is safer.
        
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            with pdfplumber.open(temp_file) as pdf:
                has_tables = False
                for i, page in enumerate(pdf.pages):
                     tables = page.extract_tables()
                     if tables:
                         has_tables = True
                         for j, table in enumerate(tables):
                             df = pd.DataFrame(table[1:], columns=table[0])
                             # Clean data?
                             sheet_name = f"Page_{i+1}_Table_{j+1}"
                             # limit sheet name len
                             df.to_excel(writer, sheet_name=sheet_name[:30], index=False)
                
                if not has_tables:
                     # fallback: Create empty with specific message
                     df = pd.DataFrame(["No tables found in PDF"], columns=["Status"])
                     df.to_excel(writer, sheet_name="Result", index=False)

        background_tasks.add_task(cleanup_file, temp_file)
        # cleanup output later
        
        return FileResponse(output_path, filename="converted_tables.xlsx", media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    except Exception as e:
        if temp_file: cleanup_file(temp_file)
        if os.path.exists(output_path): cleanup_file(output_path)
        raise HTTPException(status_code=500, detail=f"PDF to Excel failed: {str(e)}")

@router.post("/pdf-to-ppt")
async def convert_pdf_to_ppt(
    background_tasks: BackgroundTasks, 
    file: UploadFile = File(...)
):
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    
    temp_file = None
    output_filename = f"ppt_{uuid.uuid4()}.pptx"
    output_path = os.path.join(UPLOAD_DIR, output_filename)
    
    try:
        temp_file = await save_upload_file(file)
        
        # Strategy: Convert PDF pages to Images -> Slides
        # Best for visual fidelity.
        
        images = convert_from_path(temp_file)
        prs = Presentation()
        
        # Standard 16:9 
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for i, image in enumerate(images):
            # Create blank slide
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Save temp image
            img_path = os.path.join(UPLOAD_DIR, f"slide_{uuid.uuid4()}.png")
            image.save(img_path)
            
            # Add image to slide
            left = top = Inches(0)
            # Fit to slide height?
            # Ideally we check aspect ratio.
            # For simplicity, fit breadth.
            slide.shapes.add_picture(img_path, left, top, width=prs.slide_width)
            
            # Cleanup temp image immediately
            if os.path.exists(img_path):
                os.remove(img_path)

        prs.save(output_path)

        background_tasks.add_task(cleanup_file, temp_file)
        
        return FileResponse(output_path, filename="presentation.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    except Exception as e:
        if temp_file: cleanup_file(temp_file)
        if os.path.exists(output_path): cleanup_file(output_path)
        raise HTTPException(status_code=500, detail=f"PDF to PPT failed: {str(e)}")

@router.post("/pdf-to-html")
async def convert_pdf_to_html(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    temp_file = None
    output_dir = os.path.join(UPLOAD_DIR, f"pdf2html_{uuid.uuid4()}")
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        temp_file = await save_upload_file(file)
        
        # pdftohtml -s (single file) -c (complex layout) -i (ignore images? No we want them)
        # -noframes (no frameset)
        output_html = os.path.join(output_dir, "index.html")
        
        # We need to run pdftohtml
        # It comes with poppler-utils
        process = subprocess.run(
            ['pdftohtml', '-s', '-c', '-noframes', '-dataurls', temp_file, output_html],
            capture_output=True,
            text=True
        )
        
        if process.returncode != 0:
             # Fallback to simple mode if complex fails?
             raise Exception(f"pdftohtml failed: {process.stderr}")

        if not os.path.exists(output_html):
            raise Exception("Output HTML not found")

        # Read HTML content to return directly? 
        # Or return file?
        # If there are images, they might be embedded (data URI) or external.
        # With -s and -c, usually creates background images.
        # Let's inspect the directory.
        
        # Actually, for "Live Editor", receiving the HTML string is better.
        # But if there are external images, we need to serve them or embed them.
        # pdftohtml -s usually embeds? No, it might not.
        # Let's try to embed everything? 
        # -dataurls? (Available in newer poppler?)
        
        # If we return a file, the frontend needs to fetch it.
        # Let's return the HTML string if it's reasonable size, or a URL.
        # For simplicity, let's return the textual content of the HTML file, 
        # assuming basic images are handled or we zip it if complex.
        
        # For now: Return the HTML content directly.
        with open(output_html, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()

        # Cleanup
        background_tasks.add_task(shutil.rmtree, output_dir, ignore_errors=True)
        background_tasks.add_task(cleanup_file, temp_file)
        
        return Response(content=html_content, media_type="text/html")

    except Exception as e:
        if temp_file: cleanup_file(temp_file)
        shutil.rmtree(output_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"PDF to HTML failed: {str(e)}")

@router.post("/html-to-pdf")
async def convert_html_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    temp_file = None
    output_dir = os.path.join(UPLOAD_DIR, f"html2pdf_{uuid.uuid4()}")
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        temp_file = await save_upload_file(file)
        
        # Use LibreOffice to convert HTML to PDF
        process = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', temp_file, '--outdir', output_dir],
            capture_output=True,
            text=True
        )
        
        if process.returncode != 0:
             raise Exception(f"LibreOffice conversion failed: {process.stderr}")
             
        # Find PDF
        files = [f for f in os.listdir(output_dir) if f.endswith('.pdf')]
        if not files:
            raise Exception("Output PDF not found")
            
        output_pdf = os.path.join(output_dir, files[0])
        
        background_tasks.add_task(shutil.rmtree, output_dir, ignore_errors=True)
        background_tasks.add_task(cleanup_file, temp_file)
        
        return FileResponse(output_pdf, filename="converted.pdf", media_type="application/pdf")

    except Exception as e:
        if temp_file: cleanup_file(temp_file)
        shutil.rmtree(output_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"HTML to PDF failed: {str(e)}")
