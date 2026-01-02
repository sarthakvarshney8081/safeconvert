from fastapi import APIRouter, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
import os
import uuid
import shutil
from core.processor import save_upload_file, cleanup_file, UPLOAD_DIR
from pdf2docx import Converter
import pdfplumber
from pptx import Presentation
from pptx.util import Inches

router = APIRouter(prefix="/convert-from-pdf", tags=["Convert From PDF"])

@router.post("/to-word")
async def pdf_to_word(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    """
    Convert PDF to DOCX using pdf2docx.
    """
    temp_file = await save_upload_file(file)
    output_filename = os.path.splitext(os.path.basename(temp_file))[0] + ".docx"
    output_path = os.path.join(UPLOAD_DIR, output_filename)
    
    try:
        # Convert
        cv = Converter(temp_file)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
        if not os.path.exists(output_path):
             raise HTTPException(status_code=500, detail="Conversion failed")
             
        background_tasks.add_task(cleanup_file, temp_file)
        background_tasks.add_task(cleanup_file, output_path)
        
        return FileResponse(
            output_path, 
            filename=f"{os.path.splitext(file.filename)[0]}.docx",
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
    except Exception as e:
        if os.path.exists(temp_file):
            os.remove(temp_file)
        if os.path.exists(output_path):
            os.remove(output_path)
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/to-excel")
async def pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    """
    Convert PDF to Excel using pdfplumber and pandas.
    """
    import pandas as pd
    temp_file = await save_upload_file(file)
    output_filename = os.path.splitext(os.path.basename(temp_file))[0] + ".xlsx"
    output_path = os.path.join(UPLOAD_DIR, output_filename)
    
    try:
        # Simple extraction: Extract all tables
        with pdfplumber.open(temp_file) as pdf:
            all_tables = []
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    df = pd.DataFrame(table)
                    all_tables.append(df)
        
        if not all_tables:
            # Fallback: Try to just dump text or empty? 
            # If no tables, create empty excel
             pd.DataFrame(["No tables found"]).to_excel(output_path, index=False)
        else:
            # Write to multiple sheets or single? Single for now
            with pd.ExcelWriter(output_path) as writer:
                for i, df in enumerate(all_tables):
                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False, header=False)
        
        background_tasks.add_task(cleanup_file, temp_file)
        background_tasks.add_task(cleanup_file, output_path)
        
        return FileResponse(
            output_path, 
            filename=f"{os.path.splitext(file.filename)[0]}.xlsx",
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        
    except Exception as e:
        if os.path.exists(temp_file): cleanup_file(temp_file)
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/to-ppt")
async def pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    """
    Convert PDF to PPT (Image based fallback for now as direct text is hard).
    Using pdf2image -> pptx
    """
    from pdf2image import convert_from_path
    
    temp_file = await save_upload_file(file)
    output_filename = os.path.splitext(os.path.basename(temp_file))[0] + ".pptx"
    output_path = os.path.join(UPLOAD_DIR, output_filename)
    
    try:
        images = convert_from_path(temp_file)
        prs = Presentation()
        
        for img in images:
            # Save img temp
            img_path = temp_file + "_slide.jpg"
            img.save(img_path, "JPEG")
            
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to full slide
            # A4 aspect ratio usually, assuming standard PPT logic
            left = top = Inches(0)
            height = Inches(7.5) # Standard ppt height
            try:
                # auto scale
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
            except:
                 pic = slide.shapes.add_picture(img_path, left, top)
            
            if os.path.exists(img_path): os.remove(img_path)
            
        prs.save(output_path)
        
        background_tasks.add_task(cleanup_file, temp_file)
        background_tasks.add_task(cleanup_file, output_path)
        
        return FileResponse(
            output_path, 
            filename=f"{os.path.splitext(file.filename)[0]}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        if os.path.exists(temp_file): cleanup_file(temp_file)
        raise HTTPException(status_code=500, detail=str(e))
